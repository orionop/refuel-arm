#!/usr/bin/env python3
"""Generate Progress Update + Challenges tables as two PowerPoint slides."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colors ───────────────────────────────────────────────────────
NAVY       = RGBColor(0x1F, 0x35, 0x64)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY  = RGBColor(0x26, 0x26, 0x26)
GREEN      = RGBColor(0x1A, 0x7A, 0x3C)
ORANGE     = RGBColor(0xC5, 0x5A, 0x11)
BLUE       = RGBColor(0x1F, 0x6F, 0xB8)
PURPLE     = RGBColor(0x6B, 0x2F, 0xA0)
RED_DARK   = RGBColor(0xA0, 0x20, 0x20)
TEAL       = RGBColor(0x0D, 0x7C, 0x7C)

def set_cell(cell, text, bold=False, font_size=9, bg=None, fg=WHITE, align=PP_ALIGN.LEFT, italic=False):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = Pt(font_size)
    run.font.color.rgb = fg
    if bg:
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(old)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', str(bg))

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Milestones & Progress (date-wise)
# ══════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank)

# Title
txb = slide1.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.5))
tf  = txb.text_frame
p   = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Internship Progress Update — Milestones & Tasks Completed"
run.font.bold = True
run.font.size = Pt(18)
run.font.color.rgb = NAVY

milestone_data = [
    (BLUE,    "05 Jan – 16 Jan",
              "Environment Setup & Literature Survey",
              "ROS Noetic + Gazebo running with KUKA KR6 URDF; joint controllers configured; surveyed IKFast, IKFlow, MoveIt!, STOMP literature"),
    (ORANGE,  "17 Jan – 28 Jan",
              "IKFlow & cppflow Exploration → Deprecated",
              "Both approaches tested and benchmarked; deprecated due to GPU dependency, training overhead, and failure to generalise to KR6 zero-shot"),
    (BLUE,    "29 Jan – 08 Feb",
              "IK-Geo Algebraic Solver — Python Port & Validation",
              "Full port of MATLAB IK-Geo to Python; kinematic parameters (H, P) derived from URDF; up to 8 solutions/pose at ~10⁻¹⁶ m precision confirmed"),
    (GREEN,   "09 Feb – 21 Feb",
              "Custom STOMP Trajectory Optimizer Built",
              "STOMP built from scratch with smoothness matrix, joint-limit cost, and 2.5D EDT collision avoidance; cost converging from ~10,000 → ~800 in 80 iterations"),
    (TEAL,    "22 Feb – 01 Mar",
              "Full Refueling Pipeline Integrated + Topological Tests",
              "4-segment mission (HOME→YELLOW→RED→YELLOW→HOME) validated in Gazebo; Möbius strip (4π) and Pringle saddle stress-tests passed with zero elbow-flips"),
    (PURPLE,  "02 Mar – 04 Mar",
              "C-Space vs W-Space Comparison Finalised",
              "Quantitative analysis complete: C-Space requires 2 IK solves vs 60+ for W-Space; full graph suite generated and saved to output_graphs/"),
]

rows = len(milestone_data) + 1
tbl  = slide1.shapes.add_table(rows, 4, Inches(0.3), Inches(0.72), Inches(12.73), Inches(6.55)).table

tbl.columns[0].width = Inches(1.7)
tbl.columns[1].width = Inches(2.9)
tbl.columns[2].width = Inches(3.5)
tbl.columns[3].width = Inches(4.63)

headers = ["Date", "Phase", "Milestone", "Outcome / Metric"]
for ci, h in enumerate(headers):
    set_cell(tbl.cell(0, ci), h, bold=True, font_size=10, bg=NAVY, fg=WHITE, align=PP_ALIGN.CENTER)

for ri, (color, date, phase, outcome) in enumerate(milestone_data, start=1):
    row_bg = LIGHT_GREY if ri % 2 == 0 else WHITE
    set_cell(tbl.cell(ri, 0), date,    bold=True,  font_size=8.5, bg=color,   fg=WHITE)
    set_cell(tbl.cell(ri, 1), phase,   bold=True,  font_size=8.5, bg=row_bg,  fg=NAVY)
    # split outcome into milestone + metric at ';'
    parts = outcome.split(';')
    milestone_text = parts[0].strip()
    metric_text    = parts[1].strip() if len(parts) > 1 else ""
    set_cell(tbl.cell(ri, 2), milestone_text, bold=False, font_size=8,   bg=row_bg, fg=DARK_GREY)
    set_cell(tbl.cell(ri, 3), metric_text,    bold=False, font_size=8,   bg=row_bg, fg=DARK_GREY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Challenges & How They Were Overcome
# ══════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank)

txb2 = slide2.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.5))
tf2  = txb2.text_frame
p2   = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.LEFT
run2 = p2.add_run()
run2.text = "Internship Progress Update — Challenges Encountered & Overcome"
run2.font.bold = True
run2.font.size = Pt(18)
run2.font.color.rgb = NAVY

challenge_data = [
    (ORANGE,  "17 Jan – 28 Jan",
              "IKFlow & cppflow failed to generalise",
              "GPU unavailable on lab machine; both ML approaches required robot-specific training pipelines",
              "Deprecated both early; pivoted to IK-Geo — algebraic, CPU-only, zero-shot"),
    (RED_DARK,"29 Jan – 08 Feb",
              "MATLAB → Python subproblem argument mismatch",
              "IK-Geo MATLAB source used different argument ordering for SP1, SP3, SP4 than the Python library",
              "Remapped each subproblem call by cross-verifying output against FK ground truth until error reached 10⁻¹⁶ m"),
    (BLUE,    "09 Feb – 12 Feb",
              "Elbow-flips causing trajectory discontinuities",
              "IK solver returning multiple valid solutions; large joint jumps between consecutive waypoints",
              "Built solution-filtering layer ranking solutions by minimum Euclidean distance from previous config — flips eliminated"),
    (TEAL,    "13 Feb – 21 Feb",
              "STOMP rollouts violating joint limits",
              "Default Gaussian noise injection pushing trajectories into infeasible joint regions",
              "Added soft-margin joint-limit cost term penalising proximity before hard violation — trajectories stayed feasible"),
    (PURPLE,  "22 Feb – 01 Mar",
              "ROS action client timing out on long segments",
              "Fixed time_from_start spacing caused controller to reject trajectories on slower segments",
              "Implemented adaptive dt scaling proportional to joint displacement magnitude per waypoint"),
]

rows2 = len(challenge_data) + 1
tbl2  = slide2.shapes.add_table(rows2, 5, Inches(0.3), Inches(0.72), Inches(12.73), Inches(6.3)).table

tbl2.columns[0].width = Inches(1.6)
tbl2.columns[1].width = Inches(2.5)
tbl2.columns[2].width = Inches(2.8)
tbl2.columns[3].width = Inches(2.9)
tbl2.columns[4].width = Inches(2.93)

headers2 = ["Date", "Challenge", "Root Cause", "How It Was Resolved", "Lesson / Outcome"]
# merge col 3+4 conceptually — just use 5 cols
# Actually let's keep 4 cols: Date | Challenge | Root Cause | Resolution
tbl2.columns[3].width = Inches(5.83)

# redo as 4 cols
tbl2  = slide2.shapes.add_table(rows2, 4, Inches(0.3), Inches(0.72), Inches(12.73), Inches(6.3)).table
tbl2.columns[0].width = Inches(1.6)
tbl2.columns[1].width = Inches(2.8)
tbl2.columns[2].width = Inches(3.5)
tbl2.columns[3].width = Inches(4.83)

headers2 = ["Date", "Challenge", "Root Cause", "How It Was Resolved"]
for ci, h in enumerate(headers2):
    set_cell(tbl2.cell(0, ci), h, bold=True, font_size=10, bg=NAVY, fg=WHITE, align=PP_ALIGN.CENTER)

for ri, (color, date, challenge, cause, resolution) in enumerate(challenge_data, start=1):
    row_bg = LIGHT_GREY if ri % 2 == 0 else WHITE
    set_cell(tbl2.cell(ri, 0), date,        bold=True,  font_size=8.5, bg=color,  fg=WHITE)
    set_cell(tbl2.cell(ri, 1), challenge,   bold=True,  font_size=8.5, bg=row_bg, fg=NAVY)
    set_cell(tbl2.cell(ri, 2), cause,       bold=False, font_size=8,   bg=row_bg, fg=DARK_GREY)
    set_cell(tbl2.cell(ri, 3), resolution,  bold=False, font_size=8,   bg=row_bg, fg=DARK_GREY)

out = "/Users/anuragx/Desktop/Archives/projects/refuel-arm/progress_update.pptx"
prs.save(out)
print(f"Saved: {out}")
