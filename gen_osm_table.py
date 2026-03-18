#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

NAVY       = RGBColor(0x1F, 0x35, 0x64)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
DARK_GREY  = RGBColor(0x26, 0x26, 0x26)

def set_cell(cell, text, bold=False, font_size=10, bg=None, fg=WHITE, align=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.size  = Pt(font_size)
    run.font.color.rgb = fg
    if bg:
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(old)
        sf   = etree.SubElement(tcPr, qn('a:solidFill'))
        srgb = etree.SubElement(sf, qn('a:srgbClr'))
        srgb.set('val', str(bg))

data = [
    ("Objectives",
     "Solving precise 6-DOF inverse kinematics for critical refueling poses (≤ 10⁻¹³ m error) and generating smooth, joint-limit-safe trajectories between mission waypoints.\n\n"
     "Simulating and validating the full refueling sequence in a physics-based Gazebo environment while comparing C-Space vs. Workspace motion planning strategies."),
    ("Scope",
     "Executing a multi-segment refueling mission while stress-testing the IK solver through topological benchmarks (Möbius strip, Pringle saddle surface).\n\n"
     "Keeping the scope limited to simulation and algorithmic validation — excluding any hardware deployment."),
    ("Methodology",
     "Using IK-Geo for exact algebraic inverse kinematics via Paden-Kahan subproblems, and STOMP for stochastic trajectory optimization with 2.5D point cloud collision avoidance.\n\n"
     "Executing trajectories through ROS Noetic's JointTrajectoryController, with the full solver stack built in Python + NumPy using the linearSubproblemSltns library."),
]

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.5))
tf  = txb.text_frame
p   = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Objectives, Scope and Methodology"
run.font.bold = True
run.font.size = Pt(20)
run.font.color.rgb = NAVY

rows = len(data) + 1
tbl  = slide.shapes.add_table(rows, 2, Inches(0.3), Inches(0.75), Inches(12.73), Inches(6.5)).table

tbl.columns[0].width = Inches(2.2)
tbl.columns[1].width = Inches(10.53)

# Header
set_cell(tbl.cell(0, 0), "Category",    bold=True, font_size=11, bg=NAVY, fg=WHITE, align=PP_ALIGN.CENTER)
set_cell(tbl.cell(0, 1), "Description", bold=True, font_size=11, bg=NAVY, fg=WHITE, align=PP_ALIGN.CENTER)

for ri, (cat, desc) in enumerate(data, start=1):
    row_bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
    set_cell(tbl.cell(ri, 0), cat,  bold=True,  font_size=11, bg=NAVY,    fg=WHITE)
    set_cell(tbl.cell(ri, 1), desc, bold=False, font_size=10, bg=row_bg,  fg=DARK_GREY)

out = "/Users/anuragx/Desktop/objectives_scope_methodology.pptx"
prs.save(out)
print(f"Saved: {out}")
