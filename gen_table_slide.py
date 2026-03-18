#!/usr/bin/env python3
"""Generate Implementation table as a PowerPoint slide."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Colors
NAVY       = RGBColor(0x1F, 0x35, 0x64)   # header bg
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GREY  = RGBColor(0x26, 0x26, 0x26)

# Category accent colors (left border simulation via first-col bg)
CAT_COLORS = {
    "Code Organization":   RGBColor(0x1F, 0x6F, 0xB8),  # blue
    "Naming Conventions":  RGBColor(0x1F, 0x6F, 0xB8),
    "Code Documentation":  RGBColor(0x1F, 0x6F, 0xB8),
    "Readability":         RGBColor(0x21, 0x96, 0x53),  # green
    "Maintainability":     RGBColor(0x21, 0x96, 0x53),
    "Efficiency":          RGBColor(0xC5, 0x5A, 0x11),  # orange
    "Error Handling":      RGBColor(0xC5, 0x5A, 0x11),
    "Best Practices":      RGBColor(0x6B, 0x2F, 0xA0),  # purple
}

data = [
    ("Code Organization",  "Structuring the codebase into fully decoupled modules — kinematics, planning, execution, analysis",
                           "ik_geometric.py → stomp_collision.py → test_full_pipeline.py — each layer has a single responsibility"),
    ("Naming Conventions", "Using descriptive snake_case throughout — functions, variables, and constants clearly named by purpose",
                           "IK_spherical_2_parallel(), stomp_optimize(), filter_solutions(), Q_HOME, JOINT_LIMITS"),
    ("Code Documentation", "Every module opening with a docstring explaining purpose, inputs, outputs, and academic reference",
                           "STOMP module cites Kalakrishnan et al. (ICRA 2011); IK module cites MATLAB source and kinematic family"),
    ("Readability",        "Keeping functions short and single-purpose; inline comments explaining non-obvious math steps",
                           "Rodrigues rotation, subproblem mapping and pitch damping all annotated inline"),
    ("Maintainability",    "All robot-specific parameters isolated in one dictionary — changing the robot requires editing one block",
                           "KIN_KR6_R700 dict in ik_geometric.py — full pipeline adapts from this single source"),
    ("Efficiency",         "Avoiding per-waypoint IK in C-Space mode; vectorizing STOMP rollouts with NumPy batch operations",
                           "C-Space requires exactly 2 IK solves vs 60+ for W-Space; STOMP rollouts computed as 3D NumPy arrays"),
    ("Error Handling",     "Explicitly handling IK failure, joint limit violations, and ROS timeout — no silent failures",
                           "filter_solutions() returns empty array with clear print; ROS client checks result per segment"),
    ("Best Practices",     "Using argparse for CLI flags, modular __main__ guards, clean sys.path injection",
                           "All scripts runnable standalone or imported as modules without side effects"),
]

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank
slide = prs.slides.add_slide(blank_layout)

# --- Title ---
from pptx.util import Inches, Pt
txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.55))
tf  = txb.text_frame
tf.word_wrap = False
p   = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Implementation — Coding Standards & Code Quality"
run.font.bold   = True
run.font.size   = Pt(20)
run.font.color.rgb = NAVY

# --- Table dimensions ---
left   = Inches(0.3)
top    = Inches(0.8)
width  = Inches(12.73)
height = Inches(5.5)

cols = 3
rows = len(data) + 1  # +1 header

tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Column widths
tbl.columns[0].width = Inches(2.1)
tbl.columns[1].width = Inches(5.5)
tbl.columns[2].width = Inches(5.13)

# --- Helper ---
def set_cell(cell, text, bold=False, font_size=9, bg=None, fg=WHITE, align=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.size  = Pt(font_size)
    run.font.color.rgb = fg
    if bg:
        from pptx.oxml.ns import qn
        from lxml import etree
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        # remove existing solidFill
        for old in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(old)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', str(bg))

# --- Header row ---
headers = ["Category", "Aspect", "Example / Evidence"]
for ci, h in enumerate(headers):
    set_cell(tbl.cell(0, ci), h, bold=True, font_size=10, bg=NAVY, fg=WHITE, align=PP_ALIGN.CENTER)

# --- Data rows ---
for ri, (cat, aspect, evidence) in enumerate(data, start=1):
    row_bg = LIGHT_GREY if ri % 2 == 0 else WHITE
    accent = CAT_COLORS.get(cat, NAVY)

    set_cell(tbl.cell(ri, 0), cat,      bold=True,  font_size=9,  bg=accent,  fg=WHITE)
    set_cell(tbl.cell(ri, 1), aspect,   bold=False, font_size=8.5, bg=row_bg, fg=DARK_GREY)
    set_cell(tbl.cell(ri, 2), evidence, bold=False, font_size=8.5, bg=row_bg, fg=DARK_GREY)

out = "/Users/anuragx/Desktop/Archives/projects/refuel-arm/implementation_table.pptx"
prs.save(out)
print(f"Saved: {out}")
