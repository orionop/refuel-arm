#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

NAVY       = RGBColor(0x1F, 0x35, 0x64)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
DARK_GREY  = RGBColor(0x26, 0x26, 0x26)
LINK_BLUE  = RGBColor(0x1F, 0x6F, 0xB8)

def fill_cell(cell, bg):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    sf   = etree.SubElement(tcPr, qn('a:solidFill'))
    srgb = etree.SubElement(sf, qn('a:srgbClr'))
    srgb.set('val', str(bg))

def add_run(para, text, bold=False, italic=False, size=8, color=DARK_GREY, underline=False):
    run = para.add_run()
    run.text = text
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.underline = underline
    run.font.size      = Pt(size)
    run.font.color.rgb = color

def simple_cell(cell, text, bold=False, font_size=9, bg=None, fg=DARK_GREY, align=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    add_run(p, text, bold=bold, size=font_size, color=fg)
    if bg:
        fill_cell(cell, bg)

# ── Reference data (sorted latest → oldest) ──────────────────────
refs = [
    {
        "year":    "2025",
        "authors": "Bo Ai, Stephen Tian, Haochen Shi, Yixuan Wang, Tobias Pfaff, Cheston Tan, Henrik I. Christensen, Hao Su, Jiajun Wu, Yunzhu Li",
        "title":   "A review of learning-based dynamics models for robotic manipulation",
        "details": "Science Robotics, Vol. 10, eadt1497, 17 September 2025",
        "link":    "https://doi.org/10.1126/scirobotics.adt1497",
    },
    {
        "year":    "2025",
        "authors": "Alexander J. Elias, John T. Wen",
        "title":   "IK-Geo: Unified robot inverse kinematics using subproblem decomposition",
        "details": "Mechanism and Machine Theory, Vol. 209, p. 105971, ISSN 0094-114X, 2025",
        "link":    "https://doi.org/10.1016/j.mechmachtheory.2025.105971",
    },
    {
        "year":    "2024",
        "authors": "M. Yang, J. Liu",
        "title":   "Research on Six-Degree-of-Freedom Refueling Robotic Arm Positioning and Docking Based on RGB-D Visual Guidance",
        "details": "Applied Sciences, Vol. 14, No. 11, p. 4904, 2024",
        "link":    "https://doi.org/10.3390/app14114904",
    },
    {
        "year":    "2015",
        "authors": "G. P. Scott, C. G. Henshaw, I. D. Walker, B. Willimon",
        "title":   "Autonomous robotic refueling of an unmanned surface vehicle in varying sea states",
        "details": "IEEE/RSJ International Conference on Intelligent Robots and Systems (IROS), Hamburg, Germany, pp. 1664–1671, 2015",
        "link":    "https://doi.org/10.1109/IROS.2015.7353591",
    },
    {
        "year":    "2015",
        "authors": "P. Beeson, B. Ames",
        "title":   "TRAC-IK: An open-source library for improved solving of generic inverse kinematics",
        "details": "IEEE-RAS 15th International Conference on Humanoid Robots (Humanoids), Seoul, South Korea, pp. 928–935, 2015",
        "link":    "https://doi.org/10.1109/HUMANOIDS.2015.7363472",
    },
]

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Slide title ───────────────────────────────────────────────────
txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.5))
tf  = txb.text_frame
p   = tf.paragraphs[0]
run = p.add_run()
run.text = "References"
run.font.bold      = True
run.font.size      = Pt(22)
run.font.color.rgb = NAVY

# ── Table ─────────────────────────────────────────────────────────
rows = len(refs) + 1
tbl  = slide.shapes.add_table(rows, 5,
        Inches(0.3), Inches(0.72),
        Inches(12.73), Inches(6.55)).table

tbl.columns[0].width = Inches(0.7)   # Year
tbl.columns[1].width = Inches(2.4)   # Authors
tbl.columns[2].width = Inches(3.5)   # Title
tbl.columns[3].width = Inches(3.63)  # Publication Details
tbl.columns[4].width = Inches(2.5)   # Link

# Header row
for ci, h in enumerate(["Year", "Authors", "Title", "Publication Details", "Link"]):
    simple_cell(tbl.cell(0, ci), h, bold=True, font_size=10,
                bg=NAVY, fg=WHITE, align=PP_ALIGN.CENTER)

# Data rows
for ri, ref in enumerate(refs, start=1):
    row_bg = LIGHT_BLUE if ri % 2 == 0 else WHITE

    # Year — navy bg, white bold
    simple_cell(tbl.cell(ri, 0), ref["year"],
                bold=True, font_size=11, bg=NAVY, fg=WHITE, align=PP_ALIGN.CENTER)

    # Authors
    simple_cell(tbl.cell(ri, 1), ref["authors"],
                bold=False, font_size=8, bg=row_bg, fg=DARK_GREY)

    # Title — bold navy (highlighted)
    simple_cell(tbl.cell(ri, 2), ref["title"],
                bold=True, font_size=8.5, bg=row_bg, fg=NAVY)

    # Publication details
    simple_cell(tbl.cell(ri, 3), ref["details"],
                bold=False, font_size=8, bg=row_bg, fg=DARK_GREY)

    # Link — blue underline style
    cell = tbl.cell(ri, 4)
    cell.text = ""
    tf2 = cell.text_frame
    tf2.word_wrap = True
    p2  = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    add_run(p2, ref["link"], bold=False, size=7.5,
            color=LINK_BLUE, underline=True)
    fill_cell(cell, row_bg)

out = "/Users/anuragx/Desktop/references.pptx"
prs.save(out)
print(f"Saved: {out}")
